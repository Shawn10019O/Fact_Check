from __future__ import annotations
import re
from pathlib import Path
# 3rd‑party parsing libs
from pptx import Presentation
import PyPDF2
from typing import Optional,Any, cast


# ファイルからスライドまたはページのテキストを抽出
def extract_slides(path: Path) -> list[str]:
    """Return list[str] where index == slide number‑1."""
    ext = path.suffix.lower()
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pptx(path: Path) -> list[str]:
    prs = Presentation(str(path))  # mypy: Path→str
    out: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        out.append("\n".join(parts))
    return out


def _extract_pdf(path: Path) -> list[str]:
    reader = PyPDF2.PdfReader(str(path)) 
    out: list[str] = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        out.append(txt)
    return out

# テキスト前処理: 空行・改行の正規化
def sanitize(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    text = text.replace("\n", "。")
    # 句点の前に入った余計な半角スペースを除去
    text = re.sub(r" +。", "。", text)
    return text.strip("。 ")


# スライドの主題を1語で取得
def get_topic_hint(raw_slide: str, is_pptx: bool, slide_obj: Any | None = None) -> Optional[str]:
    if is_pptx and slide_obj and getattr(slide_obj.shapes, "title", None):
        title_txt = cast(str, slide_obj.shapes.title.text)
        return title_txt.strip()
    # fallback: 最初の「。」まで
    first_sentence: str = re.split(r"(?<=。)", raw_slide, maxsplit=1)[0].strip()
    print(first_sentence)
    return first_sentence[:10] 