from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from core.models import SlideResult
from core.slides import extract_slides , sanitize, get_topic_hint
from core.llm import bullets_to_paragraph, get_verdict,extract_correct_sentences
from core.search import enrich_and_filter


async def process_file(path: Path, llm_model: str) -> list[SlideResult]:
    raw_slides = extract_slides(path)
    is_pptx = path.suffix.lower() == ".pptx"
    prs = Presentation(str(path)) if is_pptx else None
    results: list[SlideResult] = []

    for idx, raw in enumerate(raw_slides, 1):
        cleaned = sanitize(raw)
        if not cleaned:
            results.append(SlideResult(idx, raw, "", "NOT_SURE", "空スライド"))
            continue
        paragraph = await bullets_to_paragraph(cleaned, llm_model)
        verdict, rationale = await get_verdict(paragraph, llm_model)

        slide_obj = prs.slides[idx - 1] if is_pptx and prs is not None else None
        hint = get_topic_hint(paragraph, is_pptx, slide_obj)

        supported_sents = await extract_correct_sentences(
            paragraph,        
            model="gpt-4o", 
            topic_hint=hint
        )

        suspicious = await enrich_and_filter(supported_sents,topic_hint=hint)
        results.append(SlideResult(idx, raw, paragraph, verdict, rationale,suspicious))
    return results