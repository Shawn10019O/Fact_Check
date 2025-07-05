import pytest, tempfile
from pathlib import Path
from pptx import Presentation
from core.process import process_file

def _dummy_pptx() -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "タイトル"
    slide.shapes.placeholders[1].text = "猫が好き。"
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return Path(tmp.name)

@pytest.mark.asyncio
async def test_process_e2e():
    res = await process_file(_dummy_pptx(), llm_model="gpt-3.5-turbo")
    assert len(res) == 1
    r = res[0]
    assert r.verdict == "SUPPORTED"
    assert r.suspicious_claims == []
